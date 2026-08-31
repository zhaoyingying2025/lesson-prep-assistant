"""教学PPT渲染引擎 - 根据LLM生成的幻灯片结构渲染为PPTX

支持多种视觉风格：
- academic: 学术简约（蓝白配色）
- cyan_ink: 青绿水墨（传统水墨）
- cute_cartoon: 清新卡通（温暖活泼）
- formal: 商务正式（深色稳重）
- minimal: 极简黑白（黑白灰）
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ============ 模板分析字体全局变量（由 export_teaching_pptx 设置） ============
_template_font_name = "Microsoft YaHei"

def _get_template_font() -> str:
    return _template_font_name

# ============ 风格配色方案 ============
STYLE_COLORS = {
    "academic": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_alt": RGBColor(0xF0, 0xF4, 0xFA),
        "primary": RGBColor(0x1A, 0x56, 0xDB),     # 蓝
        "primary_dark": RGBColor(0x0F, 0x3B, 0x9E),
        "primary_light": RGBColor(0xD6, 0xE4, 0xFD),
        "accent": RGBColor(0xE8, 0x6C, 0x00),       # 橙
        "text": RGBColor(0x1A, 0x1A, 0x2E),
        "text_light": RGBColor(0x6B, 0x72, 0x88),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "rule": RGBColor(0xE0, 0xE5, 0xEE),
        "highlight_bg": RGBColor(0xFF, 0xF3, 0xE0),
    },
    "cyan_ink": {
        "bg": RGBColor(0xFA, 0xF8, 0xF5),
        "bg_alt": RGBColor(0xF0, 0xF7, 0xF5),
        "primary": RGBColor(0x2E, 0x7D, 0x6E),     # 青绿
        "primary_dark": RGBColor(0x1C, 0x52, 0x47),
        "primary_light": RGBColor(0xD9, 0xEB, 0xE6),
        "accent": RGBColor(0xC7, 0x5C, 0x2E),       # 赭石
        "text": RGBColor(0x2C, 0x18, 0x10),          # 墨色
        "text_light": RGBColor(0x8A, 0x79, 0x68),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "rule": RGBColor(0xE8, 0xE0, 0xD8),
        "highlight_bg": RGBColor(0xFD, 0xF2, 0xE8),
    },
    "cute_cartoon": {
        "bg": RGBColor(0xFF, 0xFD, 0xFA),
        "bg_alt": RGBColor(0xFF, 0xF5, 0xF5),
        "primary": RGBColor(0xE8, 0x6B, 0x8A),     # 粉
        "primary_dark": RGBColor(0xC2, 0x4E, 0x6E),
        "primary_light": RGBColor(0xFD, 0xE0, 0xE8),
        "accent": RGBColor(0x58, 0xB3, 0x8E),       # 绿
        "text": RGBColor(0x3D, 0x2C, 0x2E),
        "text_light": RGBColor(0x9E, 0x8A, 0x8C),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "rule": RGBColor(0xEE, 0xE0, 0xE2),
        "highlight_bg": RGBColor(0xE8, 0xF5, 0xEE),
    },
    "formal": {
        "bg": RGBColor(0xF8, 0xF9, 0xFA),
        "bg_alt": RGBColor(0xEE, 0xEF, 0xF1),
        "primary": RGBColor(0x1E, 0x29, 0x3B),     # 深蓝灰
        "primary_dark": RGBColor(0x0F, 0x16, 0x24),
        "primary_light": RGBColor(0xD4, 0xD8, 0xDE),
        "accent": RGBColor(0xC0, 0x82, 0x3E),       # 金
        "text": RGBColor(0x1E, 0x1E, 0x1E),
        "text_light": RGBColor(0x6B, 0x72, 0x80),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "rule": RGBColor(0xE0, 0xE2, 0xE6),
        "highlight_bg": RGBColor(0xFE, 0xF5, 0xE7),
    },
    "minimal": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_alt": RGBColor(0xF5, 0xF5, 0xF5),
        "primary": RGBColor(0x33, 0x33, 0x33),     # 深灰
        "primary_dark": RGBColor(0x11, 0x11, 0x11),
        "primary_light": RGBColor(0xE0, 0xE0, 0xE0),
        "accent": RGBColor(0x75, 0x75, 0x75),       # 中灰
        "text": RGBColor(0x1A, 0x1A, 0x1A),
        "text_light": RGBColor(0x75, 0x75, 0x75),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "rule": RGBColor(0xE5, 0xE5, 0xE5),
        "highlight_bg": RGBColor(0xFA, 0xFA, 0xFA),
    },
}

# 16:9 标准尺寸
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 全局边距
MARGIN_L = Inches(0.7)
MARGIN_R = Inches(0.7)
MARGIN_T = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


def _get_colors(style: str) -> dict[str, RGBColor]:
    return STYLE_COLORS.get(style, STYLE_COLORS["academic"])


def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _add_rounded_rect(slide, left, top, width, height, color: RGBColor, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _add_oval(slide, left, top, width, height, color: RGBColor, line_color=None):
    """添加椭圆装饰"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _add_textbox(slide, left, top, width, height, text: str,
                 font_size=18, bold=False, color: RGBColor = None,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name=None):
    if font_name is None:
        font_name = _template_font_name
    if color is None:
        color = RGBColor(0x1A, 0x1A, 0x1A)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def _add_bullet_text(slide, left, top, width, height, items: list[str],
                     font_size=16, color: RGBColor = None, bullet_color: RGBColor = None,
                     line_spacing=1.5, font_name=None):
    if font_name is None:
        font_name = _template_font_name
    if not items:
        return
    if color is None:
        color = RGBColor(0x1A, 0x1A, 0x1A)
    if bullet_color is None:
        bullet_color = color
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # 圆点 - 使用更优雅的 ▸ 符号
        run_dot = p.add_run()
        run_dot.text = "▸  "
        run_dot.font.size = Pt(font_size)
        run_dot.font.color.rgb = bullet_color
        run_dot.font.name = font_name
        run_dot.font.bold = True
        # 正文
        run = p.add_run()
        run.text = str(item)
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def _set_bg_image(slide, image_path: str, slide_width, slide_height):
    """将图片设置为幻灯片背景（添加全屏图片并置于底层）"""
    from pptx.util import Emu
    try:
        pic = slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
    except Exception:
        pass


def _add_page_number(slide, page_num: int, total: int, colors: dict):
    """添加页码页脚条（右下角页码 + 左下角装饰线）"""
    # 底部细线
    _add_rect(slide, MARGIN_L, SLIDE_H - Inches(0.38), CONTENT_W, Inches(0.02), colors["rule"])
    # 左下角小标识
    _add_textbox(slide, MARGIN_L, SLIDE_H - Inches(0.35),
                 Inches(3), Inches(0.3),
                 "教学课件", font_size=8, color=colors["text_light"],
                 anchor=MSO_ANCHOR.MIDDLE)
    # 右下角页码
    _add_textbox(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.35),
                 Inches(1.3), Inches(0.3),
                 f"{page_num:02d} / {total:02d}", font_size=9, color=colors["text_light"],
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ============ 幻灯片渲染器 ============

def _render_cover(slide, data: dict, colors: dict, style: str):
    """封面页 - 带装饰圆形和层次感"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "教学课件")
    subtitle = data.get("subtitle", "")
    metadata = data.get("content", "")

    # 右上角装饰大圆（半透明效果用浅色模拟）
    _add_oval(slide, SLIDE_W - Inches(2.5), Inches(-0.8), Inches(3.5), Inches(3.5), colors["primary_light"])
    # 左下角装饰小圆
    _add_oval(slide, Inches(-0.5), SLIDE_H - Inches(2.0), Inches(2.5), Inches(2.5), colors["bg_alt"])

    # 顶部装饰条
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), colors["primary"])
    # 底部装饰条
    _add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), colors["primary"])

    # 中间内容卡片
    box_top = Inches(1.8)
    box_h = Inches(3.8)
    _add_rounded_rect(slide, MARGIN_L + Inches(0.3), box_top, CONTENT_W - Inches(0.6), box_h,
                      colors["bg_alt"], line_color=colors["primary_light"])

    # 卡片顶部色条
    _add_rect(slide, MARGIN_L + Inches(0.3), box_top, CONTENT_W - Inches(0.6), Inches(0.08), colors["primary"])
    # 左侧装饰竖条
    _add_rect(slide, MARGIN_L + Inches(0.3), box_top + Inches(0.08), Inches(0.06), box_h - Inches(0.08), colors["primary"])

    # 主标题
    _add_textbox(slide, MARGIN_L + Inches(1.0), box_top + Inches(0.5),
                 CONTENT_W - Inches(1.6), Inches(1.3),
                 title, font_size=38, bold=True, color=colors["primary_dark"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 副标题（章节）- 带装饰点
    if subtitle:
        _add_textbox(slide, MARGIN_L + Inches(1.0), box_top + Inches(2.0),
                     CONTENT_W - Inches(1.6), Inches(0.7),
                     f"◆  {subtitle}  ◆", font_size=20, color=colors["text_light"],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 元信息（课时/教师等）
    if metadata:
        _add_textbox(slide, MARGIN_L + Inches(1.0), box_top + Inches(2.8),
                     CONTENT_W - Inches(1.6), Inches(0.6),
                     metadata, font_size=14, color=colors["text_light"],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _render_section_header(slide, data: dict, colors: dict, style: str):
    """章节标题页 - 带装饰圆形和大色块"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "")
    subtitle = data.get("subtitle", "")

    # 右侧装饰大圆
    _add_oval(slide, SLIDE_W - Inches(3.0), Inches(0.5), Inches(4.0), Inches(4.0), colors["bg_alt"])
    # 左侧装饰小圆
    _add_oval(slide, Inches(-1.0), SLIDE_H - Inches(2.5), Inches(3.0), Inches(3.0), colors["primary_light"])

    # 全宽色块
    bar_top = Inches(2.8)
    bar_h = Inches(1.8)
    _add_rect(slide, 0, bar_top, SLIDE_W, bar_h, colors["primary"])
    # 左侧装饰条
    _add_rect(slide, 0, bar_top, Inches(0.15), bar_h, colors["primary_dark"])
    # 右侧装饰条
    _add_rect(slide, SLIDE_W - Inches(0.15), bar_top, Inches(0.15), bar_h, colors["primary_dark"])

    _add_textbox(slide, MARGIN_L + Inches(0.8), bar_top + Inches(0.15),
                 CONTENT_W - Inches(1.6), Inches(0.8),
                 title, font_size=34, bold=True, color=colors["text_white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    if subtitle:
        _add_textbox(slide, MARGIN_L + Inches(0.8), bar_top + Inches(1.0),
                     CONTENT_W - Inches(1.6), Inches(0.6),
                     f"▸  {subtitle}", font_size=18, color=colors["text_white"],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _render_content(slide, data: dict, colors: dict, style: str, page_num: int, total: int):
    """普通内容页 - 带图标标题和正文"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "")
    bullet_points = data.get("bullet_points", [])
    content = data.get("content", "")
    highlight = data.get("highlight", "")

    # 顶部装饰条
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), colors["primary"])
    # 右上角装饰圆
    _add_oval(slide, SLIDE_W - Inches(1.8), Inches(-0.6), Inches(2.2), Inches(2.2), colors["bg_alt"])

    # 标题区域 - 带左侧色条和图标背景
    _add_rounded_rect(slide, MARGIN_L, Inches(0.3), CONTENT_W, Inches(0.75), colors["bg_alt"],
                     line_color=colors["rule"])
    _add_rect(slide, MARGIN_L, Inches(0.3), Inches(0.1), Inches(0.75), colors["primary"])
    # 标题前的小方块图标
    _add_rect(slide, MARGIN_L + Inches(0.3), Inches(0.5), Inches(0.35), Inches(0.35), colors["primary"])
    _add_textbox(slide, MARGIN_L + Inches(0.8), Inches(0.3), CONTENT_W - Inches(1.1), Inches(0.75),
                 title, font_size=24, bold=True, color=colors["primary_dark"],
                 anchor=MSO_ANCHOR.MIDDLE)

    content_top = Inches(1.35)
    content_height = Inches(5.0)

    # 正文（bullet points 优先）
    if bullet_points:
        _add_bullet_text(slide, MARGIN_L + Inches(0.3), content_top,
                         CONTENT_W - Inches(0.6), content_height,
                         bullet_points, font_size=18, color=colors["text"],
                         bullet_color=colors["primary"])
    elif content:
        _add_textbox(slide, MARGIN_L + Inches(0.3), content_top,
                     CONTENT_W - Inches(0.6), content_height,
                     content, font_size=16, color=colors["text"],
                     anchor=MSO_ANCHOR.TOP)

    # 高亮内容（底部）
    if highlight:
        hl_y = content_top + content_height - Inches(0.8)
        _add_rounded_rect(slide, MARGIN_L + Inches(0.3), hl_y,
                          CONTENT_W - Inches(0.6), Inches(0.7),
                          colors["highlight_bg"], line_color=colors.get("accent", colors["primary"]))
        _add_textbox(slide, MARGIN_L + Inches(0.6), hl_y + Inches(0.05),
                     CONTENT_W - Inches(1.2), Inches(0.6),
                     f"★ {highlight}", font_size=16, bold=True,
                     color=colors.get("accent", colors["primary_dark"]),
                     anchor=MSO_ANCHOR.MIDDLE)

    _add_page_number(slide, page_num, total, colors)


def _render_two_column(slide, data: dict, colors: dict, style: str, page_num: int, total: int):
    """双栏布局页"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "")
    left_col = data.get("left_column", "")
    right_col = data.get("right_column", "")

    # 顶部装饰条
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), colors["primary"])
    _add_textbox(slide, MARGIN_L, Inches(0.3), CONTENT_W, Inches(0.7),
                 title, font_size=24, bold=True, color=colors["primary_dark"],
                 anchor=MSO_ANCHOR.MIDDLE)

    col_w = (CONTENT_W - Inches(0.3)) / 2
    col_top = Inches(1.3)
    col_h = Inches(5.2)

    # 左栏
    _add_rect(slide, MARGIN_L, col_top, col_w, col_h, colors["bg_alt"],
              line_color=colors["rule"])
    _add_rect(slide, MARGIN_L, col_top, col_w, Inches(0.04), colors["primary"])
    _add_textbox(slide, MARGIN_L + Inches(0.2), col_top + Inches(0.2),
                 col_w - Inches(0.4), col_h - Inches(0.4),
                 left_col, font_size=14, color=colors["text"])

    # 右栏
    right_x = MARGIN_L + col_w + Inches(0.3)
    _add_rect(slide, right_x, col_top, col_w, col_h, colors["bg_alt"],
              line_color=colors["rule"])
    _add_rect(slide, right_x, col_top, col_w, Inches(0.04), colors.get("accent", colors["primary"]))
    _add_textbox(slide, right_x + Inches(0.2), col_top + Inches(0.2),
                 col_w - Inches(0.4), col_h - Inches(0.4),
                 right_col, font_size=14, color=colors["text"])

    _add_page_number(slide, page_num, total, colors)


def _render_bullets(slide, data: dict, colors: dict, style: str, page_num: int, total: int):
    """纯列表页"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "")
    bullet_points = data.get("bullet_points", [])

    # 顶部装饰条
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), colors["primary"])
    _add_textbox(slide, MARGIN_L, Inches(0.3), CONTENT_W, Inches(0.7),
                 title, font_size=24, bold=True, color=colors["primary_dark"],
                 anchor=MSO_ANCHOR.MIDDLE)

    if bullet_points:
        _add_bullet_text(slide, MARGIN_L + Inches(0.5), Inches(1.4),
                         CONTENT_W - Inches(1.0), Inches(5.0),
                         bullet_points, font_size=18, color=colors["text"],
                         bullet_color=colors["primary"], line_spacing=1.8)

    _add_page_number(slide, page_num, total, colors)


def _render_table(slide, data: dict, colors: dict, style: str, page_num: int, total: int):
    """表格页"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "")
    headers = data.get("table_header", [])
    rows = data.get("table_rows", [])

    # 顶部装饰条
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), colors["primary"])
    _add_textbox(slide, MARGIN_L, Inches(0.3), CONTENT_W, Inches(0.7),
                 title, font_size=24, bold=True, color=colors["primary_dark"],
                 anchor=MSO_ANCHOR.MIDDLE)

    if not headers or not rows:
        _add_textbox(slide, MARGIN_L, Inches(2.0), CONTENT_W, Inches(1.0),
                     "暂无数据", font_size=16, color=colors["text_light"],
                     align=PP_ALIGN.CENTER)
        _add_page_number(slide, page_num, total, colors)
        return

    # 绘制表格
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 表头
    tbl_left = MARGIN_L + Inches(0.3)
    tbl_top = Inches(1.4)
    tbl_w = CONTENT_W - Inches(0.6)
    tbl_h = Inches(0.5) * min(n_rows, 10)

    shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h)
    table = shape.table

    _table_font = _template_font_name

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = colors["text_white"]
                run.font.name = _table_font
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["primary"]

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j >= n_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(13)
                    run.font.color.rgb = colors["text"]
                    run.font.name = _table_font
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["bg"] if i % 2 == 0 else colors["bg_alt"]

    _add_page_number(slide, page_num, total, colors)


def _render_quote(slide, data: dict, colors: dict, style: str, page_num: int, total: int):
    """引用/强调页"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "")
    content = data.get("content", "")
    bullet_points = data.get("bullet_points", [])

    # 顶部装饰
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), colors.get("accent", colors["primary"]))

    if title:
        _add_textbox(slide, MARGIN_L, Inches(0.3), CONTENT_W, Inches(0.7),
                     title, font_size=22, bold=True, color=colors["primary_dark"],
                     anchor=MSO_ANCHOR.MIDDLE)

    # 大引用框
    q_top = Inches(1.8)
    q_h = Inches(3.5)
    _add_rect(slide, MARGIN_L + Inches(0.5), q_top, CONTENT_W - Inches(1.0), q_h,
              colors["bg_alt"], line_color=colors.get("accent", colors["primary"]))
    # 左侧竖线
    _add_rect(slide, MARGIN_L + Inches(0.5), q_top, Inches(0.08), q_h,
              colors.get("accent", colors["primary"]))

    if bullet_points:
        display_text = "\n".join(f"● {b}" for b in bullet_points)
    else:
        display_text = content

    _add_textbox(slide, MARGIN_L + Inches(0.9), q_top + Inches(0.3),
                 CONTENT_W - Inches(1.8), q_h - Inches(0.6),
                 display_text, font_size=20, color=colors["text"],
                 anchor=MSO_ANCHOR.MIDDLE)

    _add_page_number(slide, page_num, total, colors)


def _render_activity(slide, data: dict, colors: dict, style: str, page_num: int, total: int):
    """互动活动页"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "课堂互动")
    content = data.get("content", "")
    bullet_points = data.get("bullet_points", [])

    # 顶部全宽色块
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), colors["primary"])
    _add_textbox(slide, MARGIN_L, Inches(0.15), CONTENT_W, Inches(0.7),
                 f"💬 {title}", font_size=26, bold=True, color=colors["text_white"],
                 anchor=MSO_ANCHOR.MIDDLE)

    # 活动内容
    act_top = Inches(1.4)
    act_h = Inches(5.0)

    if bullet_points:
        # 带序号的大号列表
        for i, item in enumerate(bullet_points):
            y = act_top + i * Inches(0.9)
            _add_rounded_rect(slide, MARGIN_L + Inches(0.5), y,
                              CONTENT_W - Inches(1.0), Inches(0.75),
                              colors["bg_alt"], line_color=colors["rule"])
            # 序号
            _add_rect(slide, MARGIN_L + Inches(0.5), y, Inches(0.6), Inches(0.75),
                      colors["primary"])
            _add_textbox(slide, MARGIN_L + Inches(0.5), y, Inches(0.6), Inches(0.75),
                         str(i + 1), font_size=18, bold=True, color=colors["text_white"],
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 内容
            _add_textbox(slide, MARGIN_L + Inches(1.4), y, CONTENT_W - Inches(1.9), Inches(0.75),
                         item, font_size=16, color=colors["text"],
                         anchor=MSO_ANCHOR.MIDDLE)
    elif content:
        _add_textbox(slide, MARGIN_L + Inches(0.5), act_top,
                     CONTENT_W - Inches(1.0), act_h,
                     content, font_size=18, color=colors["text"])

    _add_page_number(slide, page_num, total, colors)


def _render_summary(slide, data: dict, colors: dict, style: str, page_num: int, total: int):
    """小结页"""
    _set_bg(slide, colors["bg"])
    title = data.get("title", "本课小结")
    bullet_points = data.get("bullet_points", [])
    content = data.get("content", "")

    # 顶部全宽色块
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), colors["primary_dark"])
    _add_textbox(slide, MARGIN_L, Inches(0.15), CONTENT_W, Inches(0.7),
                 f"📋 {title}", font_size=26, bold=True, color=colors["text_white"],
                 anchor=MSO_ANCHOR.MIDDLE)

    content_top = Inches(1.5)
    content_h = Inches(5.0)

    if bullet_points:
        _add_bullet_text(slide, MARGIN_L + Inches(0.5), content_top,
                         CONTENT_W - Inches(1.0), content_h,
                         bullet_points, font_size=18, color=colors["text"],
                         bullet_color=colors["primary"], line_spacing=1.8)
    elif content:
        _add_textbox(slide, MARGIN_L + Inches(0.5), content_top,
                     CONTENT_W - Inches(1.0), content_h,
                     content, font_size=16, color=colors["text"])

    _add_page_number(slide, page_num, total, colors)


# ============ 幻灯片类型映射 ============
SLIDE_RENDERERS = {
    "cover": _render_cover,
    "section_header": _render_section_header,
    "content": _render_content,
    "two_column": _render_two_column,
    "bullets": _render_bullets,
    "table": _render_table,
    "quote": _render_quote,
    "activity": _render_activity,
    "summary": _render_summary,
}


def export_teaching_pptx(
    slide_data: dict[str, Any],
    style: str = "academic",
    template_analysis: dict | None = None,
) -> bytes:
    """将LLM生成的幻灯片结构渲染为PPTX字节流

    Args:
        slide_data: LLM生成的幻灯片数据，包含 {"slides": [...], "total_slides": N}
        style: 视觉风格标识
        template_analysis: 模板分析结果（可选），包含 layout_patterns 和 bg_image_path

    Returns:
        PPTX 文件的字节流
    """
    global _template_font_name
    _template_font_name = "Microsoft YaHei"

    colors = _get_colors(style)

    # 如果提供了模板分析结果，覆盖字体和颜色
    _template_bg_image = None
    if template_analysis:
        layout_patterns = template_analysis.get("layout_patterns", {})
        _template_bg_image = template_analysis.get("bg_image_path", "")

        if layout_patterns.get("most_common_font"):
            _template_font_name = layout_patterns["most_common_font"]

        color_palette = layout_patterns.get("color_palette", [])
        if color_palette and len(color_palette) >= 3:
            try:
                def _hex_to_rgb(h):
                    h = h.lstrip("#")
                    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
                colors["primary"] = _hex_to_rgb(color_palette[1])
                colors["text"] = _hex_to_rgb(color_palette[0])
                colors["bg"] = _hex_to_rgb(color_palette[2])
            except Exception:
                pass

    slides = slide_data.get("slides", [])
    total = len(slides)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page_num = 0
    for slide_info in slides:
        slide_type = slide_info.get("type", "content")
        renderer = SLIDE_RENDERERS.get(slide_type)

        if renderer is None:
            renderer = _render_content

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if _template_bg_image:
            _set_bg_image(slide, _template_bg_image, SLIDE_W, SLIDE_H)

        if slide_type in ("cover", "section_header"):
            renderer(slide, slide_info, colors, style)
        else:
            page_num += 1
            renderer(slide, slide_info, colors, style, page_num, total)

        notes_text = slide_info.get("speaker_notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ============ HTML预览生成 ============

STYLE_COLORS_CSS = {
    "academic": {"bg": "#FFFFFF", "bg_alt": "#F0F4FA", "primary": "#1A56DB", "primary_dark": "#0F3B9E", "primary_light": "#D6E4FD", "accent": "#E86C00", "text": "#1A1A2E", "text_light": "#6B7288", "text_white": "#FFFFFF", "rule": "#E0E5EE", "highlight_bg": "#FFF3E0"},
    "cyan_ink": {"bg": "#FAF8F5", "bg_alt": "#F0F7F5", "primary": "#2E7D6E", "primary_dark": "#1C5247", "primary_light": "#D9EBE6", "accent": "#C75C2E", "text": "#2C1810", "text_light": "#8A7968", "text_white": "#FFFFFF", "rule": "#E8E0D8", "highlight_bg": "#FDF2E8"},
    "cute_cartoon": {"bg": "#FFFDFA", "bg_alt": "#FFF5F5", "primary": "#E86B8A", "primary_dark": "#C24E6E", "primary_light": "#FDE0E8", "accent": "#58B38E", "text": "#3D2C2E", "text_light": "#9E8A8C", "text_white": "#FFFFFF", "rule": "#EEE0E2", "highlight_bg": "#E8F5EE"},
    "formal": {"bg": "#F8F9FA", "bg_alt": "#EEEFF1", "primary": "#1E293B", "primary_dark": "#0F1624", "primary_light": "#D4D8DE", "accent": "#C0823E", "text": "#1E1E1E", "text_light": "#6B7280", "text_white": "#FFFFFF", "rule": "#E0E2E6", "highlight_bg": "#FEF5E7"},
    "minimal": {"bg": "#FFFFFF", "bg_alt": "#F5F5F5", "primary": "#333333", "primary_dark": "#111111", "primary_light": "#E0E0E0", "accent": "#757575", "text": "#1A1A1A", "text_light": "#757575", "text_white": "#FFFFFF", "rule": "#E5E5E5", "highlight_bg": "#FAFAFA"},
}


def _escape_html(text: str) -> str:
    """简单的HTML转义"""
    return (str(text)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#39;"))


def _render_slide_html(slide_data: dict, colors: dict, style: str, page_num: int, total: int) -> str:
    """将单个幻灯片数据渲染为HTML预览"""
    slide_type = slide_data.get("type", "content")
    title = slide_data.get("title", "")
    content = slide_data.get("content", "")
    bullet_points = slide_data.get("bullet_points", [])
    highlight = slide_data.get("highlight", "")
    subtitle = slide_data.get("subtitle", "")
    left_col = slide_data.get("left_column", "")
    right_col = slide_data.get("right_column", "")
    headers = slide_data.get("table_header", [])
    rows = slide_data.get("table_rows", [])
    speaker_notes = slide_data.get("speaker_notes", "")

    c = colors
    is_cover_or_section = slide_type in ("cover", "section_header")

    # 标题栏
    def _title_bar(extra_cls=""):
        return f"""
        <div class="flex items-center gap-2 mb-3 {extra_cls}">
            <div class="w-1 h-7 rounded" style="background:{c['primary']}"></div>
            <div class="font-bold" style="font-size:20px;color:{c['primary_dark']}">{_escape_html(title)}</div>
        </div>"""

    def _bullet_html(items):
        if not items:
            return ""
        bullets = "".join(f'<li class="mb-1.5 leading-relaxed" style="color:{c["text"]}"><span style="color:{c["primary"]};font-weight:bold;margin-right:6px">▸</span>{_escape_html(item)}</li>' for item in items)
        return f'<ul class="list-none p-0 m-0">{bullets}</ul>'

    def _page_num():
        if is_cover_or_section:
            return ""
        return f"""
        <div class="flex items-center justify-between px-4 py-1" style="border-top:1px solid {c['rule']}">
            <span style="font-size:9px;color:{c['text_light']}">教学课件</span>
            <span style="font-size:9px;color:{c['text_light']}">{page_num:02d} / {total:02d}</span>
        </div>"""

    if slide_type == "cover":
        return f"""
        <div class="slide-preview" style="background:{c['bg']};border-radius:8px;overflow:hidden;position:relative;min-height:280px;border:1px solid {c['rule']}">
            <div style="position:absolute;top:-20px;right:-20px;width:120px;height:120px;border-radius:50%;background:{c['primary_light']};opacity:0.6"></div>
            <div style="position:absolute;bottom:-20px;left:-20px;width:80px;height:80px;border-radius:50%;background:{c['bg_alt']};opacity:0.6"></div>
            <div style="height:3px;background:{c['primary']}"></div>
            <div class="flex flex-col items-center justify-center px-6 py-8" style="min-height:240px">
                <div style="font-size:28px;font-weight:bold;color:{c['primary_dark']};text-align:center;margin-bottom:8px">{_escape_html(title)}</div>
                {f'<div style="font-size:16px;color:{c["text_light"]};text-align:center;margin-bottom:6px">◆ {_escape_html(subtitle)} ◆</div>' if subtitle else ''}
                {f'<div style="font-size:12px;color:{c["text_light"]};text-align:center">{_escape_html(content)}</div>' if content else ''}
            </div>
            <div style="height:3px;background:{c['primary']}"></div>
        </div>"""

    if slide_type == "section_header":
        return f"""
        <div class="slide-preview" style="background:{c['bg']};border-radius:8px;overflow:hidden;position:relative;min-height:200px;border:1px solid {c['rule']}">
            <div style="position:absolute;top:10px;right:-30px;width:100px;height:100px;border-radius:50%;background:{c['bg_alt']};opacity:0.6"></div>
            <div style="position:absolute;bottom:-20px;left:-30px;width:80px;height:80px;border-radius:50%;background:{c['primary_light']};opacity:0.6"></div>
            <div class="flex flex-col items-center justify-center" style="background:{c['primary']};margin:60px 0;padding:20px;min-height:100px">
                <div style="font-size:24px;font-weight:bold;color:{c['text_white']};text-align:center">{_escape_html(title)}</div>
                {f'<div style="font-size:14px;color:{c["text_white"]};text-align:center;margin-top:6px">▸ {_escape_html(subtitle)}</div>' if subtitle else ''}
            </div>
        </div>"""

    # 通用内容页
    inner = ""
    if slide_type == "two_column":
        inner = f"""
        <div class="flex gap-2 mt-1">
            <div class="flex-1 p-3 rounded" style="background:{c['bg_alt']};border:1px solid {c['rule']}">
                <div style="height:2px;background:{c['primary']};margin-bottom:6px"></div>
                <div style="font-size:13px;color:{c['text']};white-space:pre-wrap">{_escape_html(left_col)}</div>
            </div>
            <div class="flex-1 p-3 rounded" style="background:{c['bg_alt']};border:1px solid {c['rule']}">
                <div style="height:2px;background:{c['accent']};margin-bottom:6px"></div>
                <div style="font-size:13px;color:{c['text']};white-space:pre-wrap">{_escape_html(right_col)}</div>
            </div>
        </div>"""
    elif slide_type == "table" and headers and rows:
        tbl_rows = ""
        for i, row in enumerate(rows):
            bg = c['bg'] if i % 2 == 0 else c['bg_alt']
            cells = "".join(f'<td class="px-2 py-1.5 text-center" style="font-size:12px;color:{c["text"]};background:{bg}">{_escape_html(str(v))}</td>' for v in row)
            tbl_rows += f"<tr>{cells}</tr>"
        hdr = "".join(f'<th class="px-2 py-1.5 text-center font-bold" style="font-size:13px;color:{c["text_white"]};background:{c["primary"]}">{_escape_html(h)}</th>' for h in headers)
        inner = f"""
        <div class="mt-1 overflow-x-auto">
            <table class="w-full border-collapse rounded overflow-hidden" style="border:1px solid {c['rule']}">
                <thead><tr>{hdr}</tr></thead>
                <tbody>{tbl_rows}</tbody>
            </table>
        </div>"""
    elif slide_type == "quote":
        inner = f"""
        <div class="mt-1 p-4 rounded" style="background:{c['bg_alt']};border-left:4px solid {c['accent']};border:1px solid {c['rule']};border-left-width:4px">
            {"".join(f'<div class="mb-1" style="font-size:14px;color:{c["text"]}">● {_escape_html(b)}</div>' for b in bullet_points) if bullet_points else f'<div style="font-size:14px;color:{c["text"]};white-space:pre-wrap">{_escape_html(content)}</div>'}
        </div>"""
    elif slide_type == "activity":
        inner = f"""
        <div class="flex flex-col gap-1.5 mt-1">
            {''.join(f'<div class="flex items-center gap-2 p-2 rounded" style="background:{c["bg_alt"]};border:1px solid {c["rule"]}"><div class="w-5 h-5 flex items-center justify-center rounded font-bold text-xs" style="background:{c["primary"]};color:{c["text_white"]}">{i+1}</div><div style="font-size:13px;color:{c["text"]}">{_escape_html(b)}</div></div>' for i, b in enumerate(bullet_points)) if bullet_points else f'<div style="font-size:13px;color:{c["text"]};white-space:pre-wrap">{_escape_html(content)}</div>'}
        </div>"""
    elif slide_type == "summary":
        inner = _bullet_html(bullet_points) if bullet_points else (f'<div style="font-size:13px;color:{c["text"]};white-space:pre-wrap">{_escape_html(content)}</div>' if content else "")
    else:
        # content / bullets / default
        inner = _bullet_html(bullet_points) if bullet_points else (f'<div style="font-size:13px;color:{c["text"]};white-space:pre-wrap">{_escape_html(content)}</div>' if content else "")

    highlight_html = ""
    if highlight:
        highlight_html = f"""
        <div class="mt-2 p-2 rounded text-sm font-bold" style="background:{c['highlight_bg']};border:1px solid {c['accent']};color:{c['accent']}">
            ★ {_escape_html(highlight)}
        </div>"""

    notes_html = ""
    if speaker_notes:
        notes_html = f"""<div class="mt-1 text-xs" style="color:{c['text_light']}">🎤 {_escape_html(speaker_notes)}</div>"""

    return f"""
    <div class="slide-preview" style="background:{c['bg']};border-radius:8px;overflow:hidden;border:1px solid {c['rule']}">
        <div style="height:2px;background:{c['primary']}"></div>
        <div class="p-3">
            {_title_bar()}
            {inner}
            {highlight_html}
            {notes_html}
        </div>
        {_page_num()}
    </div>"""


def generate_ppt_preview_html(
    slide_data: dict[str, Any],
    style: str = "academic",
) -> str:
    """将幻灯片数据生成为HTML预览字符串

    Args:
        slide_data: LLM生成的幻灯片数据，包含 {"slides": [...], "total_slides": N}
        style: 视觉风格标识

    Returns:
        完整的HTML预览字符串
    """
    colors = STYLE_COLORS_CSS.get(style, STYLE_COLORS_CSS["academic"])
    slides = slide_data.get("slides", [])
    total = len(slides)

    page_num = 0
    slides_html = []
    for i, slide_info in enumerate(slides):
        slide_type = slide_info.get("type", "content")
        if slide_type in ("cover", "section_header"):
            slides_html.append(_render_slide_html(slide_info, colors, style, 0, total))
        else:
            page_num += 1
            slides_html.append(_render_slide_html(slide_info, colors, style, page_num, total))

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>PPT预览</title>
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{ font-family:-apple-system,BlinkMacSystemFont,"Microsoft YaHei","PingFang SC",sans-serif; background:#f0f0f0; padding:20px; }}
.slide-preview {{ max-width:820px; margin:0 auto 16px auto; box-shadow:0 2px 8px rgba(0,0,0,0.1); }}
.slide-preview:hover {{ box-shadow:0 4px 16px rgba(0,0,0,0.15); }}
</style>
</head>
<body>
{"".join(slides_html)}
</body>
</html>"""