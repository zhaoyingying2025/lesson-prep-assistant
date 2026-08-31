"""教学PPT生成 Agent：将教案转化为课堂使用的教学幻灯片"""
from __future__ import annotations

import json
import io
from collections import Counter
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Emu

from ..core.llm import get_llm
from ..core.prompts import PPT_SYSTEM, PPT_USER_TEMPLATE
from ..core.prompt_loader import inject_domain_context
from ..models.schemas import LessonPlan
from ..storage.file_store import save_upload


# 风格配置
PPT_STYLES = {
    "academic": {
        "label": "学术简约",
        "desc": "干净清晰，蓝白配色，适合理工科",
    },
    "cyan_ink": {
        "label": "青绿水墨",
        "desc": "传统水墨风格，淡雅配色，适合文科/艺术",
    },
    "cute_cartoon": {
        "label": "清新卡通",
        "desc": "温暖活泼，圆角设计，适合低年级",
    },
    "formal": {
        "label": "商务正式",
        "desc": "深色稳重，适合MBA/管理类课程",
    },
    "minimal": {
        "label": "极简黑白",
        "desc": "纯黑白灰，极度简洁，适合任何场合",
    },
}

DENSITY_DESC = {
    "detailed": "内容详细丰富，每页3-5个要点，适合自学或复习课",
    "moderate": "内容适中，每页2-3个要点，适合常规课堂教学",
    "concise": "内容精炼简洁，每页1-2个要点，适合快速讲授或复习",
}

IMAGE_DESC = {
    "none": "无图，纯文字排版，用表格、对比、流程图等文字化视觉元素",
    "icons": "使用图标装饰，用文字符号（如★●▶◆）辅助视觉层级",
    "rich": "使用丰富的视觉表达，如分栏、表格、对比框、引用块等布局变化",
}


async def generate_ppt_content(
    plan: LessonPlan,
    knowledge_points: list[dict],
    style: str = "academic",
    content_density: str = "moderate",
    image_style: str = "icons",
    style_custom: str = "",
    textbook_context: str = "",
    subject: Optional[str] = None,
    template_analysis: Optional[dict] = None,
) -> dict[str, Any]:
    """调用LLM生成教学PPT的幻灯片结构

    Args:
        plan: 完整教案
        knowledge_points: 知识点列表
        style: PPT视觉风格
        content_density: 内容密度
        image_style: 视觉元素风格
        style_custom: 用户自定义风格描述
        textbook_context: 教材原文参考，用于知识点解释时引用教材原文
        subject: 学科标识(如 math/chinese/english/physics 等), 用于注入学科领域规则
        template_analysis: 模板分析结果（可选），包含 layout_patterns 等排版信息

    Returns:
        {"style_used": str, "total_slides": int, "slides": [...]}
    """
    llm = get_llm()

    params = {
        "style": PPT_STYLES.get(style, PPT_STYLES["academic"])["label"],
        "density": content_density,
        "content_density_desc": DENSITY_DESC.get(content_density, DENSITY_DESC["moderate"]),
        "image_style_desc": IMAGE_DESC.get(image_style, IMAGE_DESC["icons"]),
        "style_custom": style_custom,
    }

    # 如果提供了模板分析结果，注入排版信息到 LLM 上下文
    if template_analysis:
        lp = template_analysis.get("layout_patterns", {})
        if lp:
            template_info = (
                f"\n【上传模板的排版参考】\n"
                f"- 模板尺寸：{lp.get('slide_width_cm', '?')}cm × {lp.get('slide_height_cm', '?')}cm\n"
                f"- 推荐字体：{lp.get('most_common_font', '微软雅黑')}（常用字号：{lp.get('most_common_font_size_pt', 18)}pt）\n"
                f"- 配色方案：{', '.join(lp.get('color_palette', [])[:5])}\n"
                f"- 正文区域（左/上边距）：{lp.get('margin_left_cm', '?')}cm / {lp.get('body_area_cm', {}).get('top', '?')}cm\n"
                f"- 标题位置：顶部 {lp.get('title_position_cm', {}).get('top', '?')}cm\n"
                f"- 是否使用项目符号：{'是' if lp.get('has_bullets') else '否'}\n"
            )
            if template_analysis.get("bg_image_path"):
                template_info += f"- 模板包含背景图，将在生成时应用\n"
            params["template_reference"] = template_info

    # 课程序列化
    plan_dict = plan.model_dump()
    # 将 stages 等对象转为 dict
    plan_dict["stages"] = [s.model_dump() for s in plan.stages]

    user_prompt = PPT_USER_TEMPLATE.format(
        course_name=plan.course_name,
        chapter=plan.chapter,
        total_minutes=plan.total_minutes,
        lesson_json=json.dumps(plan_dict, ensure_ascii=False, indent=2),
        knowledge_json=json.dumps(knowledge_points, ensure_ascii=False, indent=2),
        textbook_context=textbook_context or "（暂无教材原文参考，请基于知识点内容生成）",
        params_json=json.dumps(params, ensure_ascii=False, indent=2),
        style=style,
        density=content_density,
        content_density_desc=params["content_density_desc"],
        image_style_desc=params["image_style_desc"],
    )

    # 注入学科领域规则 (借鉴 ai-teaching-ppt 的多槽位注入)
    system_prompt = inject_domain_context(PPT_SYSTEM, subject)

    data = await llm.chat_json(system_prompt, user_prompt, temperature=0.8)

    # 验证结构
    if "slides" not in data or not isinstance(data["slides"], list):
        raise ValueError("PPT生成失败：LLM返回结构异常")

    return data


async def analyze_pptx_template(pptx_path: str) -> dict:
    """分析PPTX模板文件，提取排版习惯和背景图

    Returns:
        {"layout_patterns": {...}, "bg_image_path": ""}
    """
    prs = PptxPresentation(pptx_path)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    all_fonts = []
    all_font_sizes = []
    all_colors = []
    text_box_info = []
    has_bullets = False
    bullet_count = 0
    shape_type_counts = Counter()
    bg_image_extracted = None

    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    for slide in prs.slides:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            fill_type_name = str(fill.type)
        else:
            fill_type_name = "none"

        if bg_image_extracted is None:
            try:
                if fill.type is not None:
                    from pptx.enum.dml import MSO_THEME_COLOR
                    if hasattr(fill, "fill") and fill.fill is not None:
                        inner = fill.fill
                    else:
                        inner = fill
                    if hasattr(inner, "user_photo"):
                        pass
            except Exception:
                pass

        slide_layout = slide.slide_layout
        placeholder_count = len(slide_layout.placeholders)

        for shape in slide.shapes:
            shape_type_counts[shape.shape_type] += 1

            if shape.has_text_frame:
                tf = shape.text_frame
                box_info = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "word_wrap": tf.word_wrap,
                    "paragraphs": [],
                }
                for para in tf.paragraphs:
                    para_info = {"alignment": str(para.alignment) if para.alignment else "left", "runs": []}
                    if para.level > 0:
                        has_bullets = True
                        bullet_count += 1
                    for run in para.runs:
                        font = run.font
                        font_name = font.name
                        font_size = font.size
                        bold = font.bold
                        italic = font.italic
                        color = None
                        if font.color and font.color.type is not None:
                            try:
                                color = str(font.color.rgb)
                            except Exception:
                                color = str(font.color.theme_color) if font.color.theme_color else None
                            if color:
                                all_colors.append(color)
                        if font_name:
                            all_fonts.append(font_name)
                        if font_size:
                            all_font_sizes.append(font_size.pt)
                        para_info["runs"].append({
                            "font": font_name,
                            "size": font_size.pt if font_size else None,
                            "bold": bool(bold),
                            "italic": bool(italic),
                        })
                    box_info["paragraphs"].append(para_info)
                text_box_info.append(box_info)

    font_counter = Counter(all_fonts)
    size_counter = Counter(round(s, 1) for s in all_font_sizes)
    color_counter = Counter(all_colors)

    most_common_font = font_counter.most_common(1)[0][0] if font_counter else "微软雅黑"
    most_common_size = size_counter.most_common(1)[0][0] if size_counter else 18.0
    top_colors = [c for c, _ in color_counter.most_common(5)]

    title_positions = []
    body_positions = []
    for info in text_box_info:
        if info["top"] < slide_height // 2 and info["height"] < slide_height // 3:
            title_positions.append({"top": info["top"], "left": info["left"]})
        else:
            body_positions.append({"top": info["top"], "left": info["left"]})

    avg_title_top = (
        round(sum(p["top"] for p in title_positions) / len(title_positions) / 914400 * 2.54, 1)
        if title_positions else 2.5
    )
    avg_body_left = (
        round(sum(p["left"] for p in body_positions) / len(body_positions) / 914400 * 2.54, 1)
        if body_positions else 1.5
    )
    avg_body_top = (
        round(sum(p["top"] for p in body_positions) / len(body_positions) / 914400 * 2.54, 1)
        if body_positions else 4.0
    )

    margin_left = round(min((info["left"] for info in text_box_info), default=914400) / 914400 * 2.54, 1)
    margin_right = round(
        (slide_width - max((info["left"] + info["width"] for info in text_box_info), default=0)) / 914400 * 2.54, 1
    )

    layout_patterns = {
        "slide_width_cm": round(slide_width / 914400 * 2.54, 1),
        "slide_height_cm": round(slide_height / 914400 * 2.54, 1),
        "slide_count": len(prs.slides),
        "total_text_boxes": len(text_box_info),
        "avg_placeholder_count": round(placeholder_count, 1),
        "most_common_font": most_common_font,
        "most_common_font_size_pt": most_common_size,
        "font_variety": len(font_counter),
        "top_fonts": [{"name": n, "count": c} for n, c in font_counter.most_common(5)],
        "font_sizes_used": sorted(set(round(s, 1) for s in all_font_sizes)),
        "color_palette": top_colors,
        "color_variety": len(color_counter),
        "has_bullets": has_bullets,
        "bullet_ratio": round(bullet_count / max(len(text_box_info), 1), 2),
        "title_position_cm": {"top": avg_title_top, "left": 0.0},
        "body_area_cm": {"left": avg_body_left, "top": avg_body_top},
        "margin_left_cm": margin_left,
        "margin_right_cm": margin_right,
        "shape_type_breakdown": {str(k): v for k, v in shape_type_counts.most_common(10)},
        "background_fill_type": fill_type_name,
    }

    bg_image_path = ""
    for slide in prs.slides:
        try:
            bg = slide.background
            if bg.fill.type is not None:
                from pptx.oxml.ns import qn
                bg_elem = bg._element
                blip_fill = bg_elem.find(qn("p:blipFill"))
                if blip_fill is not None:
                    blip = blip_fill.find(qn("a:blip"))
                    if blip is not None:
                        embed_id = blip.get(qn("r:embed"))
                        if embed_id:
                            image_part = prs.part.related_parts.get(embed_id)
                            if image_part:
                                image_bytes = image_part.blob
                                ext = "png"
                                content_type = getattr(image_part, "content_type", "") or ""
                                if "jpeg" in content_type or "jpg" in content_type:
                                    ext = "jpg"
                                elif "gif" in content_type:
                                    ext = "gif"
                                elif "bmp" in content_type:
                                    ext = "bmp"
                                bg_image_path = str(
                                    await save_upload(image_bytes, f"bg_template_{Path(pptx_path).stem}.{ext}", 0)
                                )
                                break
        except Exception:
            continue

    return {
        "layout_patterns": layout_patterns,
        "bg_image_path": bg_image_path,
    }